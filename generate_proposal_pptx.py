import pptx
from pptx.util import Inches

# Create a PowerPoint Presentation
prs = pptx.Presentation()

# Title Slide
slide_title = prs.slides.add_slide(prs.slide_layouts[0])
slide_title.shapes.title.text = "PKF Proposal & Code Opening Automation"
slide_title.shapes.placeholders[1].text = "Generated on 2026-01-28"

# Architecture Diagram Slide
slide_architecture = prs.slides.add_slide(prs.slide_layouts[5])
slide_architecture.shapes.title.text = "Architecture Diagram"

# Add architecture diagram image (this is a placeholder)
# Actual image path should be added here
slide_architecture.shapes.add_picture('architecture_diagram.png', Inches(0.5), Inches(1.5), height=Inches(3))

# Process Flows Slide
slide_process_flows = prs.slides.add_slide(prs.slide_layouts[5])
slide_process_flows.shapes.title.text = "Process Flows"

# Add process flow diagram image (this is a placeholder)
# Actual image path should be added here
slide_process_flows.shapes.add_picture('process_flows.png', Inches(0.5), Inches(1.5), height=Inches(3))

# Implementation Plan Slide
slide_impl_plan = prs.slides.add_slide(prs.slide_layouts[5])
slide_impl_plan.shapes.title.text = "Implementation Plan"

# Add implementation plan content
content = "1. Phase 1: Research\n2. Phase 2: Development\n3. Phase 3: Testing\n4. Phase 4: Deployment"

text_box = slide_impl_plan.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(5))
text_frame = text_box.text_frame
text_frame.text = content

# Save the presentation
prs.save('PKF_Proposal.pptx')
