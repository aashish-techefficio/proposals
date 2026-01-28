import pptx
from pptx.util import Inches

# Create a PowerPoint presentation object
presentation = pptx.Presentation()

# Define a function to add a slide with a title and content

def add_slide(title, content):
    slide_layout = presentation.slide_layouts[1]  # Use layout with title and content
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

# Slide 1: Title
add_slide("PKF Proposal & Code Opening Automation", "Generated on: 2026-01-28 15:56:04 UTC")

# Slide 2: Architecture
add_slide("Architecture", "This slide illustrates the architecture of the PKF Proposal & Code Opening Automation solution.")

# Slide 3: Scope
add_slide("Scope", "The scope of the project includes automation of proposal generation and code opening processes.")

# Slide 4: Implementation Plan
add_slide("Implementation Plan", "1. Requirement Gathering\n2. Design Phase\n3. Development Phase\n4. Testing Phase\n5. Deployment")

# Slide 5: Proposal Workflow
add_slide("Proposal Workflow", "Step 1: Idea Generation\nStep 2: Drafting Proposal\nStep 3: Review and Approval\nStep 4: Final Submission")

# Slide 6: Code Opening Workflow
add_slide("Code Opening Workflow", "1. Code Review\n2. Integration\n3. Final Testing\n4. Code Release")

# Slide 7: Highlights
add_slide("Highlights", "- Efficient Proposal Generation\n- Streamlined Code Opening Process\n- Improved Collaboration")

# Slide 8: Next Steps
add_slide("Next Steps", "- Review this proposal\n- Schedule a follow-up meeting\n- Start the implementation phase")

# Save the presentation
presentation.save('PKF_Proposal_Generator.pptx')
