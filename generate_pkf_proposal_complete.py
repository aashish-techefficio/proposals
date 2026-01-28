import pptx
from pptx.util import Inches

# Create a presentation object
prs = pptx.Presentation()

# Title Slide
slide1 = prs.slides.add_slide(prs.slide_layouts[0]) # layout for title slide
title = slide1.shapes.title
subtitle = slide1.placeholders[1]
title.text = "PKF Proposal & Code Opening Automation"
subtitle.text = "A Comprehensive Solution for Automation"

# Architecture Diagram
slide2 = prs.slides.add_slide(prs.slide_layouts[5]) # blank layout
slide2.shapes.title.text = 'Architecture Diagram'
# Assuming you have an image path for the architecture diagram
# slide2.shapes.add_picture('architecture_diagram.png', Inches(1), Inches(1), width=Inches(8), height=Inches(4))

# Scope
slide3 = prs.slides.add_slide(prs.slide_layouts[1]) # title and content layout
slide3.shapes.title.text = 'Scope'
content = slide3.placeholders[1]
content.text = 'This proposal covers the automation of PKF proposal generation and code opening.'

# Implementation Plan
slide4 = prs.slides.add_slide(prs.slide_layouts[1])
slide4.shapes.title.text = 'Implementation Plan'
content = slide4.placeholders[1]
content.text = 'Step 1: Requirement Analysis\nStep 2: Development\nStep 3: Testing\nStep 4: Deployment'

# Timelines
slide5 = prs.slides.add_slide(prs.slide_layouts[1])
slide5.shapes.title.text = 'Timelines'
content = slide5.placeholders[1]
content.text = 'Total Duration: 12 weeks\nWeek 1-2: Requirement Analysis\nWeek 3-6: Development\nWeek 7-8: Testing\nWeek 9-12: Deployment'

# Process Flows
slide6 = prs.slides.add_slide(prs.slide_layouts[5])
slide6.shapes.title.text = 'Process Flows'
# Assuming you have a picture of a process flow
# slide6.shapes.add_picture('process_flow.png', Inches(1), Inches(1), width=Inches(8), height=Inches(4))

# Conclusion
slide7 = prs.slides.add_slide(prs.slide_layouts[1])
slide7.shapes.title.text = 'Conclusion'
content = slide7.placeholders[1]
content.text = 'This proposal provides a detailed plan for automating PKF proposal generation and code opening, enhancing efficiency and accuracy.'

# Save the presentation
prs.save('PKF_Proposal_Automation.pptx')
