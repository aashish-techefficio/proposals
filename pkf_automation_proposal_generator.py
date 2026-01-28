from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor

# Function to set font styles
def set_font(para, font_name, font_size, white_font=False):
    run = para.add_run()
    run.font.name = font_name
    run.font.size = Pt(font_size)
    if white_font:
        run.font.color.rgb = RGBColor(255, 255, 255)
    return run

# Function to create a slide
def create_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title  # Set slide title
    textbox = slide.shapes.placeholders[1]
    textbox.text = content  # Set slide content

# Create presentation
presentation = Presentation()

# Title Slide
slide_layout = presentation.slide_layouts[0]  # Title Slide Layout
slide = presentation.slides.add_slide(slide_layout)
slide.shapes.title.text = "PKF Proposal & Code Opening Automation"
slide.shapes.placeholders[1].text = "Generated on: 2026-01-28"

# Content Slides
slides_content = [
    ("Introduction", "This is the introduction of the PKF Proposal & Code Opening Automation solution which encompasses various aspects of automation..."),
    ("Objectives", "The main objectives of this proposal include improving efficiency, reducing errors, and ensuring quality code..."),
    ("Proposed Solution", "The proposed solution involves several steps, including...")
]  

for title, content in slides_content:
    create_slide(presentation, title, content)

# Adding a table slide
slide_layout = presentation.slide_layouts[5]  # Title Only Slide
slide = presentation.slides.add_slide(slide_layout)

# Let's create a table
table = slide.shapes.add_table(rows=3, cols=3, left=Inches(2), top=Inches(1), width=Inches(6), height=Inches(2.5)).table

# Assigning cell values
table.cell(0, 0).text = "Feature"
table.cell(0, 1).text = "Description"
table.cell(0, 2).text = "Benefits"

table.cell(1, 0).text = "Automation"
table.cell(1, 1).text = "Automating code openings and proposals"
table.cell(1, 2).text = "Increased efficiency and accuracy"

table.cell(2, 0).text = "Quality Assurance"
table.cell(2, 1).text = "Ensuring predefined quality standards"
table.cell(2, 2).text = "Consistent delivery of high-quality output"

# Save the presentation
presentation.save("PKF_Proposal_Automation.pptx")
