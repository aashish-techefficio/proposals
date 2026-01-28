from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
presentation = Presentation()

# Add a slide for the architecture diagram
slide_layout = presentation.slide_layouts[5]  # Use a blank slide layout
slide = presentation.slides.add_slide(slide_layout)

# Define elements of the diagram
shapes = []

# Users
users = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(1), Inches(2), Inches(1))
users.text = 'Users'
shapes.append(users)

# SharePoint Online
sharepoint = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4), Inches(1), Inches(2), Inches(1))
sharepoint.text = 'SharePoint Online'
shapes.append(sharepoint)

# Power Automate
power_automate = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(3), Inches(2), Inches(1))
power_automate.text = 'Power Automate'
shapes.append(power_automate)

# Reporting layers
reporting = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4), Inches(3), Inches(2), Inches(1))
reporting.text = 'Reporting Layer'
shapes.append(reporting)

# Connect the elements with arrows
for i, shape_start in enumerate(shapes):
    for shape_end in shapes[i+1:]:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, shape_start.left + shape_start.width/2, shape_start.top + shape_start.height/2,
                                           shape_end.left + shape_end.width/2, shape_end.top + shape_end.height/2)

# Save the presentation
presentation.save('Architecture_Diagram.pptx')

print('PowerPoint presentation created successfully!')