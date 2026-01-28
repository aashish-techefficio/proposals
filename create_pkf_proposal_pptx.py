import os
from pptx import Presentation
from pptx.util import Inches

class PKFProposalPresentation:
    def __init__(self):
        self.presentation = Presentation()

    def add_title_slide(self, title, subtitle):
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[0])
        title_placeholder = slide.shapes.title
        subtitle_placeholder = slide.placeholders[1]
        title_placeholder.text = title
        subtitle_placeholder.text = subtitle

    def add_architecture_diagram(self, image_path):
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[5])
        slide.shapes.add_picture(image_path, Inches(0.5), Inches(0.5), width=Inches(8.5))

    def add_process_flows(self, flow_data):
        # Placeholder for process flows content
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[5])
        title = slide.shapes.title
        title.text = 'Process Flows'
        # Add flowchart as needed

    def add_scope(self, scope_text):
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = 'Scope'
        content.text = scope_text

    def add_implementation_plan(self, plan_text):
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = 'Implementation Plan'
        content.text = plan_text

    def add_timelines(self, timeline_data):
        # Placeholder for timelines content
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[5])
        title = slide.shapes.title
        title.text = 'Timelines'
        # Add timeline visual as needed

    def add_expected_outcomes(self, outcomes_text):
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = 'Expected Outcomes'
        content.text = outcomes_text

    def add_team_roles(self, team_data):
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = 'Team Roles'
        content.text = team_data

    def add_conclusion(self, conclusion_text):
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = 'Conclusion'
        content.text = conclusion_text

    def save(self, file_name):
        self.presentation.save(file_name)

if __name__ == '__main__':
    presentation = PKFProposalPresentation()
    presentation.add_title_slide('PKF Proposal & Code Opening Automation', 'Generated on 2026-01-28')
    # Add actual content by calling other methods here
    # Example: presentation.add_scope('This project will cover...')
    presentation.save('PKF_Proposal_Presentation.pptx')
